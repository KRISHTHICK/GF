# ===============================
# Required Packages (pip install)
# ===============================
# pip install streamlit
# pip install pdfplumber
# pip install pytesseract
# pip install transformers
# pip install torch
# pip install numpy
# pip install pandas
# pip install python-pptx
# pip install pillow

import streamlit as st
import pdfplumber
import pytesseract
from PIL import Image
from transformers import pipeline
from pptx import Presentation
import json
import numpy as np
import pandas as pd
import os
import tempfile

# Load high-accuracy NER model
@st.cache_resource
def load_ner_model():
    return pipeline("ner", model="Jean-Baptiste/roberta-large-ner-english", grouped_entities=True)

# Make output JSON serializable
def make_serializable(obj):
    if isinstance(obj, dict):
        return {k: make_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [make_serializable(i) for i in obj]
    elif isinstance(obj, (np.float32, np.float64)):
        return float(obj)
    elif isinstance(obj, (np.int32, np.int64)):
        return int(obj)
    return obj

# Extract text from PDF
def extract_text_from_pdf(file):
    try:
        with pdfplumber.open(file) as pdf:
            return "\n".join([page.extract_text() or "" for page in pdf.pages])
    except Exception as e:
        st.error(f"PDF error: {e}")
        return ""

# Extract text from image using OCR
def extract_text_from_image(file):
    try:
        image = Image.open(file)
        return pytesseract.image_to_string(image)
    except Exception as e:
        st.error(f"Image error: {e}")
        return ""

# Extract text from PPTX
def extract_text_from_pptx(file):
    try:
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        st.error(f"PPTX error: {e}")
        return ""

# Perform NER and get top N entities
def extract_entities(text, model, top_n=20):
    try:
        entities = model(text)
        entities = make_serializable(entities)
        df = pd.DataFrame(entities)
        df = df.sort_values(by="score", ascending=False).drop_duplicates(subset=["word"]).head(top_n)
        return df.to_dict(orient="records")
    except Exception as e:
        st.error(f"NER error: {e}")
        return []

# Main Streamlit app
def main():
    st.title("Multi-Format Entity Extractor (PDF, Image, PPT)")

    uploaded_file = st.file_uploader("Upload a file (PDF, PNG, JPEG, PPTX)", type=["pdf", "png", "jpg", "jpeg", "pptx"])

    if uploaded_file:
        st.info(f"Uploaded: {uploaded_file.name}")

        file_ext = os.path.splitext(uploaded_file.name)[-1].lower()
        text = ""

        with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp_file:
            tmp_file.write(uploaded_file.read())
            tmp_path = tmp_file.name

        if file_ext == ".pdf":
            text = extract_text_from_pdf(tmp_path)
        elif file_ext in [".png", ".jpg", ".jpeg"]:
            text = extract_text_from_image(tmp_path)
        elif file_ext == ".pptx":
            text = extract_text_from_pptx(tmp_path)

        if not text.strip():
            st.error("No text extracted from the file.")
            return

        st.subheader("Extracted Text (Preview)")
        st.text(text[:1000])

        ner_model = load_ner_model()
        top_n = st.slider("How many top entities to extract?", 10, 50, 20)

        entities = extract_entities(text, ner_model, top_n=top_n)

        if entities:
            st.subheader("Extracted Entities (JSON)")
            st.json(entities)

            st.download_button(
                label="Download JSON",
                data=json.dumps(entities, indent=4),
                file_name="extracted_entities.json",
                mime="application/json"
            )
        else:
            st.warning("No entities found.")

if __name__ == "__main__":
    main()
